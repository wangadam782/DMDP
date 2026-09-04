#!/usr/bin/env python
from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


DEFAULT_COMPARISON_DIR = Path("outputs/comparisons/four_methods_with_method4_residual_warmup_100k")
DEFAULT_OUTPUT = Path("outputs/presentations/method4_latest_comparison_results.pptx")

METRICS = [
    ("mean_return", "Return", "higher"),
    ("success_rate", "Success", "higher"),
    ("average_cost", "Cost", "lower"),
    ("state_w2", "State W2", "lower"),
    ("final_state_w2", "Final W2", "lower"),
    ("tail_risk", "Tail Risk", "lower"),
    ("unsafe_occupancy", "Unsafe", "lower"),
    ("dispersion_error", "Dispersion", "lower"),
]


def load_rows(path: Path) -> list[dict[str, Any]]:
    return json.loads(path.read_text(encoding="utf-8"))


def set_wide_layout(prs: Presentation) -> None:
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)


def add_footer(slide, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.45), Inches(7.05), Inches(12.4), Inches(0.28))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(105, 105, 105)


def add_title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "方法四最新实验结果对比"
    slide.placeholders[1].text = "Tail Warmup 与 Residual Warmup 的 100k 训练评估；对比方法一、二、三基线"
    add_footer(slide, "Environment: Safety-Gymnasium MultiGoal, 10-episode evaluation")


def add_bullet_slide(prs: Presentation, title: str, bullets: list[str], footer: str | None = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(21)
    if footer:
        add_footer(slide, footer)


def fmt(value: Any) -> str:
    if value is None:
        return ""
    return f"{float(value):.3f}"


def add_table_slide(prs: Presentation, rows: list[dict[str, Any]]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "统一评估结果表"

    table = slide.shapes.add_table(
        len(rows) + 1,
        len(METRICS) + 1,
        Inches(0.25),
        Inches(1.15),
        Inches(12.85),
        Inches(4.55),
    ).table
    table.cell(0, 0).text = "Method"
    for col, (_, title, _) in enumerate(METRICS, start=1):
        table.cell(0, col).text = title

    for r, row in enumerate(rows, start=1):
        table.cell(r, 0).text = row["method"].replace("Method ", "M")
        for c, (key, _, _) in enumerate(METRICS, start=1):
            table.cell(r, c).text = fmt(row.get(key))

    for row_idx in range(len(rows) + 1):
        for col_idx in range(len(METRICS) + 1):
            cell = table.cell(row_idx, col_idx)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(9 if col_idx else 8)
                paragraph.alignment = PP_ALIGN.CENTER
                if row_idx == 0:
                    paragraph.font.bold = True

    note = slide.shapes.add_textbox(Inches(0.55), Inches(5.95), Inches(12.1), Inches(0.7))
    tf = note.text_frame
    tf.text = "结论：Tail Warmup 是当前方法四主结果；Residual Warmup 的确定性评估退回方法三水平，暂不作为主结果。"
    tf.paragraphs[0].font.size = Pt(17)
    tf.paragraphs[0].font.bold = True
    add_footer(slide, "Higher is better: Return, Success. Lower is better: Cost, W2, Tail Risk, Unsafe, Dispersion.")


def add_two_figures_slide(
    prs: Presentation,
    title: str,
    left_image: Path,
    right_image: Path,
    bullets: list[str],
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    slide.shapes.add_picture(str(left_image), Inches(0.45), Inches(1.1), width=Inches(6.1))
    slide.shapes.add_picture(str(right_image), Inches(6.85), Inches(1.1), width=Inches(6.1))
    box = slide.shapes.add_textbox(Inches(0.7), Inches(5.95), Inches(12.0), Inches(0.85))
    tf = box.text_frame
    tf.clear()
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(15)
    add_footer(slide, f"Figures: {left_image.parent}")


def add_needed_figures_slide(prs: Presentation) -> None:
    bullets = [
        "主文建议放 4 张核心柱状图：Return、Average Cost、Tail Risk、State W2。",
        "补充材料建议放：Success Rate、Unsafe Occupancy、Final State W2、Dispersion Error。",
        "若论文强调理论方法四，需要再补 1 张训练曲线图：recent episode return / cost / tail risk 随 step 变化。",
        "若解释 Residual Warmup 失败，需要补 1 张 residual action norm 或 residual mean norm 曲线，说明确定性残差没有形成稳定控制。",
        "最终论文叙述建议：方法四 Tail Warmup 同时改善任务收益和 cost，但 W2 指标仍需进一步优化。",
    ]
    add_bullet_slide(prs, "论文需要的结果插图", bullets)


def build_presentation(rows: list[dict[str, Any]], comparison_dir: Path, output: Path) -> None:
    figure_dir = comparison_dir / "figures"
    prs = Presentation()
    set_wide_layout(prs)

    add_title_slide(prs)
    add_bullet_slide(
        prs,
        "本轮实验设置",
        [
            "对比对象：方法一 State-MAPPO、方法二 Distributional MAPPO、方法三 DMDP-MAPPO、方法四 Online / Tail Warmup / Residual Warmup。",
            "方法四 Tail Warmup：cost penalty 从 0.10 warmup 到 0.12，同时逐步打开 smooth tail risk 与 Lyapunov penalty。",
            "方法四 Residual Warmup：以方法三 actor 为 frozen base，方法四学习 residual policy。",
            "评估口径：统一 checkpoint，10 episodes，指标包括任务回报、成功率、cost、W2、tail risk、unsafe occupancy。",
        ],
        "Result source: eval_summary.json",
    )
    add_table_slide(prs, rows)
    add_two_figures_slide(
        prs,
        "任务收益与安全代价",
        figure_dir / "mean_return.png",
        figure_dir / "average_cost.png",
        [
            "Tail Warmup 的 Return 最高，且 Cost 低于方法一、二、三，是当前最强方法四版本。",
            "Online 原版虽然 Return 较高，但 Cost 明显过高，安全性不足。",
        ],
    )
    add_two_figures_slide(
        prs,
        "尾部风险与不安全占用",
        figure_dir / "tail_risk.png",
        figure_dir / "unsafe_occupancy.png",
        [
            "方法三在 Tail Risk / Unsafe 上仍是最安全的基线。",
            "Tail Warmup 相比 Online 原版大幅降低尾部风险，但仍略高于方法三。",
        ],
    )
    add_two_figures_slide(
        prs,
        "状态分布距离",
        figure_dir / "state_w2.png",
        figure_dir / "final_state_w2.png",
        [
            "Online 原版 Final W2 最低，但以高 cost 和高 unsafe 为代价。",
            "Tail Warmup 的 W2 不占优，说明下一轮应继续加强分布约束，而不是继续加大 tail penalty。",
        ],
    )
    add_bullet_slide(
        prs,
        "当前结论与下一轮方向",
        [
            "主结果建议采用 Method 4 Tail Warmup：Return=2.824，Success=0.80，Cost=137.9。",
            "Residual Warmup 不建议作为主结果：确定性评估与方法三完全一致，说明 residual 均值控制没有稳定生效。",
            "下一轮方法四应保留 Tail Warmup 的 reward schedule，并调小但延长 W2 penalty：例如 distribution_penalty=0.003-0.008，warmup_steps=50k。",
            "同时记录 omega_delta、Lyapunov violation、tail risk time series，用于证明方法四符合论文中的在线分布反馈控制思想。",
        ],
    )
    add_needed_figures_slide(prs)

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--comparison-dir", default=str(DEFAULT_COMPARISON_DIR))
    parser.add_argument("--output", default=str(DEFAULT_OUTPUT))
    args = parser.parse_args()

    comparison_dir = Path(args.comparison_dir)
    rows = load_rows(comparison_dir / "eval_summary.json")
    build_presentation(rows, comparison_dir, Path(args.output))
    print(f"saved presentation to {args.output}")


if __name__ == "__main__":
    main()
