#!/usr/bin/env python
from __future__ import annotations

import argparse
import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


RUN_DIR_DEFAULT = Path("outputs/runs/mappo_method1_multigoal_100k")


def load_json(path: Path) -> dict:
    return json.loads(path.read_text(encoding="utf-8"))


def set_wide_layout(prs: Presentation) -> None:
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_bullet_slide(prs: Presentation, title: str, bullets: list[str], note: str | None = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(20)
    if note:
        add_footer(slide, note)


def add_footer(slide, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.5), Inches(6.9), Inches(12.2), Inches(0.35))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(95, 95, 95)


def add_placeholder(slide, left: float, top: float, width: float, height: float, title: str, body: str) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 247, 250)
    line = shape.line
    line.color.rgb = RGBColor(160, 160, 160)

    tf = shape.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    p1.text = title
    p1.font.bold = True
    p1.font.size = Pt(18)
    p1.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = body
    p2.font.size = Pt(13)
    p2.alignment = PP_ALIGN.CENTER


def add_training_setup_slide(prs: Presentation, train: dict) -> None:
    config = train["config"]
    adapter = train["adapter"]
    bullets = [
        f"环境：{adapter['actual_env_id']}，2-agent Safe Multi-Agent MultiGoal。",
        "策略形式：a_t^i = pi_theta(o_t^i)，仅输入局部观测，不输入 omega_t。",
        f"训练步数：{train['summary']['total_steps']}；完成 episode 数：{train['summary']['num_completed_episodes']}。",
        f"PPO 设置：rollout_steps={config['rollout_steps']}，num_epochs={config['num_epochs']}，clip_ratio={config['clip_ratio']}。",
        f"折扣与优势：gamma={config['gamma']}，gae_lambda={config['gae_lambda']}。",
        f"网络：hidden_sizes={config['network']['hidden_sizes']}，activation={config['network']['activation']}。",
        f"安全项：基础 cost penalty = {config['cost_penalty']}，不包含 W2 或 tail-risk penalty。",
    ]
    add_bullet_slide(prs, "方法一训练设置", bullets)


def add_results_table_slide(prs: Presentation, mappo_eval: dict, random_eval: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "方法一结果摘要"

    rows = [
        ("Mean Return", f"{mappo_eval['metrics']['mean_return']:.3f}", f"{random_eval['metrics']['mean_return']:.3f}"),
        ("Success Rate", f"{mappo_eval['metrics']['success_rate']:.2f}", f"{random_eval['metrics']['success_rate']:.2f}"),
        ("Average Cost", f"{mappo_eval['metrics']['average_cost']:.1f}", f"{random_eval['metrics']['average_cost']:.1f}"),
        ("Return CVaR 0.1", f"{mappo_eval['metrics']['return_cvar_0.1']:.3f}", f"{random_eval['metrics']['return_cvar_0.1']:.3f}"),
        ("State W2", f"{mappo_eval['metrics']['state_w2']:.3f}", f"{random_eval['metrics']['state_w2']:.3f}"),
        ("Tail Risk", f"{mappo_eval['metrics']['tail_risk']:.2f}", f"{random_eval['metrics']['tail_risk']:.2f}"),
        ("Unsafe Occupancy", f"{mappo_eval['metrics']['unsafe_occupancy']:.2f}", f"{random_eval['metrics']['unsafe_occupancy']:.2f}"),
    ]

    table = slide.shapes.add_table(len(rows) + 1, 3, Inches(0.8), Inches(1.4), Inches(7.0), Inches(4.8)).table
    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "MAPPO"
    table.cell(0, 2).text = "Random"
    for i, (name, mappo_value, random_value) in enumerate(rows, start=1):
        table.cell(i, 0).text = name
        table.cell(i, 1).text = mappo_value
        table.cell(i, 2).text = random_value

    bullets = [
        "当前结果说明：方法一已经明显优于 random，对任务回报和 cost 都有提升。",
        "但 Tail Risk 和 Unsafe Occupancy 仍然偏高，不能说明其状态分布已经安全。",
        "因此它适合作为方法三 DMDP-MAPPO 的主要对照基线。",
    ]
    box = slide.shapes.add_textbox(Inches(8.2), Inches(1.5), Inches(4.3), Inches(3.8))
    tf = box.text_frame
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(18)


def add_image_plan_slide(prs: Presentation, title: str, items: list[tuple[str, str, str]], footer: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    positions = [
        (0.7, 1.4, 5.8, 2.0),
        (6.8, 1.4, 5.8, 2.0),
        (0.7, 3.8, 5.8, 2.0),
        (6.8, 3.8, 5.8, 2.0),
    ]
    for (label, filename, message), (left, top, width, height) in zip(items, positions):
        add_placeholder(
            slide,
            left,
            top,
            width,
            height,
            f"建议插图：{label}",
            f"{filename}\n{message}",
        )
    add_footer(slide, footer)


def build_presentation(train: dict, mappo_eval: dict, random_eval: dict, output: Path, figure_dir: Path) -> None:
    prs = Presentation()
    set_wide_layout(prs)

    add_title_slide(
        prs,
        "方法一：State-Feedback MAPPO",
        "用途：作为 DMDP-MAPPO 的对照基线；不使用状态分布反馈，只依赖局部观测。",
    )

    add_bullet_slide(
        prs,
        "方法一定位",
        [
            "方法一是标准 State-Feedback MAPPO，用来回答：只靠局部状态反馈时，能把任务做到什么程度。",
            "策略输入：a_t^i = pi_theta(o_t^i)。",
            "训练目标：最大化 reward，并可使用基础 cost penalty：r_base = r_env - lambda_c * cost。",
            "它不使用 DMDP 的 omega_t = [mu_t, sigma_t]，也不直接约束系统状态分布 rho_t。",
            "因此它适合与方法三比较：在回报接近时，状态分布安全性是否更差。",
        ],
    )

    add_bullet_slide(
        prs,
        "与方法三的核心区别",
        [
            "方法一关注：局部观测 -> 动作，主要优化任务回报。",
            "方法三关注：局部观测 + 状态分布参数 omega_t -> 动作，显式调节 rho_t 接近 rho_star。",
            "因此，方法一即使任务回报不错，也不代表 State W2、Tail Risk、Unsafe Occupancy 会低。",
            "本 PPT 的目标不是证明方法一更好，而是说明它是必要且合理的对照组。",
        ],
    )

    add_training_setup_slide(prs, train)

    add_image_plan_slide(
        prs,
        "训练过程建议插图",
        [
            ("训练回报曲线", "train_recent_episode_return.png", "用于说明方法一已经学到有效策略。"),
            ("训练 cost 曲线", "train_recent_episode_cost.png", "用于说明回报提升是否伴随安全代价。"),
            ("value loss 曲线", "train_value_loss.png", "用于说明 critic 是否稳定。"),
            ("entropy 曲线", "train_entropy.png", "用于说明探索是否过早塌缩。"),
        ],
        f"图像目录建议：{figure_dir}",
    )

    add_results_table_slide(prs, mappo_eval, random_eval)

    add_image_plan_slide(
        prs,
        "结果展示建议插图",
        [
            ("任务表现", "mean_return.png", "展示方法一相对 random 的任务收益提升。"),
            ("成功率", "success_rate.png", "说明方法一是否真正完成任务，而不只是偶然拿到回报。"),
            ("安全代价", "average_cost.png", "说明方法一是否靠更高风险换取任务收益。"),
            ("状态分布差异", "state_w2.png", "说明方法一与目标安全分布 rho_star 的距离。"),
        ],
        f"图像目录建议：{figure_dir}",
    )

    add_bullet_slide(
        prs,
        "当前可得结论",
        [
            f"训练到 {train['summary']['total_steps']} steps 后，方法一已完成 {train['summary']['num_completed_episodes']} 个 episode，可视为有效 baseline。",
            f"最近训练窗口的 recent_episode_return = {train['summary']['recent_episode_return']:.3f}，说明策略已不再是随机行为。",
            f"最终 10-episode 评估中，MAPPO mean_return = {mappo_eval['metrics']['mean_return']:.3f}，高于 random 的 {random_eval['metrics']['mean_return']:.3f}。",
            f"MAPPO average_cost = {mappo_eval['metrics']['average_cost']:.1f}，显著低于 random 的 {random_eval['metrics']['average_cost']:.1f}。",
            f"MAPPO state_w2 = {mappo_eval['metrics']['state_w2']:.3f}，低于 random 的 {random_eval['metrics']['state_w2']:.3f}，但 Tail Risk 与 Unsafe Occupancy 仍不理想。",
        ],
    )

    add_bullet_slide(
        prs,
        "与方法三需要对比的问题",
        [
            "在相近 Mean Return 和 Return CVaR 下，DMDP-MAPPO 是否能进一步降低 State W2？",
            "在相近 Success Rate 下，DMDP-MAPPO 是否能进一步降低 Tail Risk 和 Unsafe Occupancy？",
            "方法一是否主要依靠局部策略改进，而没有真正改善系统状态分布？",
            "若方法三引入 omega_t 后安全指标更优，则可以支持“状态分布控制”比“纯状态反馈”更有效的假设。",
        ],
    )

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--run-dir", default=str(RUN_DIR_DEFAULT))
    parser.add_argument("--output", default="outputs/presentations/method1_mappo_overview.pptx")
    args = parser.parse_args()

    run_dir = Path(args.run_dir)
    train = load_json(run_dir / "metrics" / "mappo_train.json")
    mappo_eval = load_json(run_dir / "eval" / "mappo_eval.json")
    random_eval = load_json(run_dir / "eval" / "random_eval.json")
    figure_dir = run_dir / "figures"
    build_presentation(train, mappo_eval, random_eval, Path(args.output), figure_dir)
    print(f"saved presentation to {args.output}")


if __name__ == "__main__":
    main()
