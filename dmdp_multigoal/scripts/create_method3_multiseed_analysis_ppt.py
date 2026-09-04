#!/usr/bin/env python
from __future__ import annotations

import argparse
import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


DEFAULT_SUMMARY = Path("outputs/comparisons/three_methods_3seeds/eval_summary.json")
DEFAULT_OUT = Path("outputs/presentations/method3_multiseed_analysis.pptx")
DEFAULT_FIG_DIR = Path("outputs/comparisons/three_methods_3seeds/figures")


def load_json(path: Path) -> dict:
    return json.loads(path.read_text(encoding="utf-8"))


def set_wide_layout(prs: Presentation) -> None:
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)


def add_footer(slide, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.45), Inches(6.9), Inches(12.3), Inches(0.3))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(95, 95, 95)


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_bullet_slide(prs: Presentation, title: str, bullets: list[str], footer: str | None = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(19)
    if footer:
        add_footer(slide, footer)


def add_placeholder(slide, left: float, top: float, width: float, height: float, title: str, body: str) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(245, 247, 250)
    shape.line.color.rgb = RGBColor(160, 160, 160)
    tf = shape.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    p1.text = title
    p1.font.bold = True
    p1.font.size = Pt(17)
    p1.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = body
    p2.font.size = Pt(12)
    p2.alignment = PP_ALIGN.CENTER


def metric_line(summary: dict, method: str, key: str) -> str:
    cell = summary[method]["metrics"][key]
    return f"{cell['mean']:.4f} ± {cell['std']:.4f}"


def add_stats_table_slide(prs: Presentation, summary: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "三方法统计对比"

    columns = [
        ("Method", 2.35),
        ("Return", 1.45),
        ("Success", 1.25),
        ("Cost", 1.35),
        ("CVaR0.1", 1.45),
        ("State W2", 1.35),
        ("Tail", 1.1),
        ("Unsafe", 1.2),
    ]
    rows = [
        ("State-Feedback MAPPO", "State-Feedback MAPPO"),
        ("Distributional MAPPO", "Distributional MAPPO"),
        ("DMDP-MAPPO", "DMDP-MAPPO"),
    ]
    metric_keys = [
        "mean_return",
        "success_rate",
        "average_cost",
        "return_cvar_0.1",
        "state_w2",
        "tail_risk",
        "unsafe_occupancy",
    ]

    left = Inches(0.3)
    top = Inches(1.45)
    width = Inches(sum(col_width for _, col_width in columns))
    height = Inches(2.5)
    table = slide.shapes.add_table(len(rows) + 1, len(columns), left, top, width, height).table

    for idx, (_, col_width) in enumerate(columns):
        table.columns[idx].width = Inches(col_width)

    for col_idx, (label, _) in enumerate(columns):
        cell = table.cell(0, col_idx)
        cell.text = label
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(12)
            p.alignment = PP_ALIGN.CENTER

    for row_idx, (display_name, summary_name) in enumerate(rows, start=1):
        cell = table.cell(row_idx, 0)
        cell.text = display_name
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
        for col_idx, metric_key in enumerate(metric_keys, start=1):
            metric = summary[summary_name]["metrics"][metric_key]
            value = f"{metric['mean']:.3f} ± {metric['std']:.3f}"
            metric_cell = table.cell(row_idx, col_idx)
            metric_cell.text = value
            for p in metric_cell.text_frame.paragraphs:
                p.font.size = Pt(10.5)
                p.alignment = PP_ALIGN.CENTER

    box = slide.shapes.add_textbox(Inches(0.5), Inches(4.35), Inches(12.2), Inches(1.6))
    tf = box.text_frame
    bullets = [
        "建议保留这页作为汇报中的核心数据页，先讲 Return / Success / Cost，再讲 State W2 / Tail / Unsafe。",
        "如果版面紧张，可把表格保留，把详细原因分析放到下一页。",
    ]
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(17)
    add_footer(slide, "数据来源：outputs/comparisons/three_methods_3seeds/eval_summary.json")


def add_plan_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "下一步验证计划"

    left_box = slide.shapes.add_textbox(Inches(0.55), Inches(1.35), Inches(6.05), Inches(4.8))
    tf_left = left_box.text_frame
    left_bullets = [
        "1. 提高 omega_t 质量：用多个并行环境共同估计 rho_hat_t，而不是每步只用 2 个 agent 样本。",
        "2. 对 mu_t、sigma_t 做时间平滑，降低当前状态分布反馈的抖动。",
        "3. 将 tail 约束从硬阈值统计改为平滑距离惩罚，提高可学习性。",
    ]
    for idx, bullet in enumerate(left_bullets):
        p = tf_left.paragraphs[0] if idx == 0 else tf_left.add_paragraph()
        p.text = bullet
        p.font.size = Pt(19)

    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.35), Inches(5.95), Inches(4.8))
    tf_right = right_box.text_frame
    right_bullets = [
        "4. 固定方法一和方法二，不再修改基线，只继续优化方法三。",
        "5. 统一训练预算、网络规模、评估协议，并使用 3 到 5 个 seeds 做公平比较。",
        "6. 理论验证标准：在 Mean Return 和 Success Rate 接近的前提下，DMDP 的 State W2、Tail Risk、Unsafe Occupancy 更低。",
    ]
    for idx, bullet in enumerate(right_bullets):
        p = tf_right.paragraphs[0] if idx == 0 else tf_right.add_paragraph()
        p.text = bullet
        p.font.size = Pt(19)

    add_footer(slide, "这页用于说明：下一步不是堆更多实验，而是先补齐 DMDP 的关键理论部件，再做公平多 seed 验证。")


def build_presentation(summary: dict, out_path: Path, fig_dir: Path) -> None:
    prs = Presentation()
    set_wide_layout(prs)

    add_title_slide(
        prs,
        "方法三多 Seed 分析",
        "主题：为什么当前 DMDP-MAPPO 在 3 个 seed 下没有稳定优于 Distributional MAPPO。",
    )

    add_bullet_slide(
        prs,
        "结果结论",
        [
            f"方法二 Mean Return: {metric_line(summary, 'Distributional MAPPO', 'mean_return')}；方法三: {metric_line(summary, 'DMDP-MAPPO', 'mean_return')}。",
            f"方法二 Success Rate: {metric_line(summary, 'Distributional MAPPO', 'success_rate')}；方法三: {metric_line(summary, 'DMDP-MAPPO', 'success_rate')}。",
            f"方法三唯一稳定优势是 State W2: {metric_line(summary, 'DMDP-MAPPO', 'state_w2')}，略低于方法二的 {metric_line(summary, 'Distributional MAPPO', 'state_w2')}。",
            f"但方法三的 Tail Risk: {metric_line(summary, 'DMDP-MAPPO', 'tail_risk')}，Unsafe Occupancy: {metric_line(summary, 'DMDP-MAPPO', 'unsafe_occupancy')}，没有稳定优于方法二。",
            "因此当前版本的方法三只能证明“更接近目标分布的均值/方差”，还不能证明“尾部安全更强”。",
        ],
        footer="建议插图：左侧插入 mean_return.png；右侧插入 state_w2.png；底部可补 tail_risk.png。",
    )

    add_stats_table_slide(prs, summary)

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "原因分析"
    add_placeholder(
        slide,
        0.6,
        1.4,
        6.0,
        2.0,
        "建议插图：任务表现对比",
        "文件：mean_return.png\n放在左上。\n用于说明方法三在 return / success 上没有稳定压过方法二。",
    )
    add_placeholder(
        slide,
        6.9,
        1.4,
        5.8,
        2.0,
        "建议插图：状态分布对比",
        "文件：state_w2.png\n放在右上。\n用于说明方法三的主要收益集中在 W2 上。",
    )
    box = slide.shapes.add_textbox(Inches(0.7), Inches(3.9), Inches(12.0), Inches(2.4))
    tf = box.text_frame
    bullets = [
        "1. 当前 omega_t 每步只用单环境的 2 个 agent 样本估计，分布反馈噪声仍然偏大。",
        "2. 当前训练主要优化 W2，tail_penalty=0.0，因此方法三更像 distribution matching，而不是 tail-aware safety control。",
        "3. 对角高斯只保留均值和方差，忽略危险状态的相关结构，所以 W2 下降不等于 tail 风险下降。",
        "4. hand-crafted rho_star 只表达平均安全，不表达尾部安全，这会限制方法三的上限。",
    ]
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(18)
    add_footer(slide, f"图像目录：{fig_dir}")

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "当前建议"
    add_placeholder(
        slide,
        0.6,
        1.5,
        5.9,
        1.8,
        "建议插图：tail_risk 对比",
        "文件：tail_risk.png\n放在左上。\n说明当前方法三没有稳定压低 tail 指标。",
    )
    add_placeholder(
        slide,
        6.9,
        1.5,
        5.8,
        1.8,
        "建议插图：unsafe_occupancy 对比",
        "文件：unsafe_occupancy.png\n放在右上。\n说明尾部占用没有形成稳定优势。",
    )
    box = slide.shapes.add_textbox(Inches(0.75), Inches(3.7), Inches(12.0), Inches(2.3))
    tf = box.text_frame
    steps = [
        "1. 保留现有 hand-crafted target，不再继续使用当前两版 data-driven target。",
        "2. 优先把 omega_t 估计改成真正的多并行环境统计，而不是单环境 2-agent 估计。",
        "3. 如果继续做 tail 控制，改成平滑距离惩罚，不直接用当前硬阈值统计。",
        "4. 在不动方法一和方法二的前提下，只继续改方法三，再做多 seed 复验。",
    ]
    for idx, step in enumerate(steps):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = step
        p.font.size = Pt(18)
    add_footer(slide, "这页用于交代决策：短期不换 target，先提高方法三分布估计质量。")

    add_plan_slide(prs)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--summary", default=str(DEFAULT_SUMMARY))
    parser.add_argument("--fig-dir", default=str(DEFAULT_FIG_DIR))
    parser.add_argument("--output", default=str(DEFAULT_OUT))
    args = parser.parse_args()

    summary = load_json(Path(args.summary))
    build_presentation(summary, Path(args.output), Path(args.fig_dir))
    print(f"saved presentation to {args.output}")


if __name__ == "__main__":
    main()
