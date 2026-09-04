#!/usr/bin/env python
from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]


def rel(path: str | Path) -> str:
    return str((ROOT / path).resolve())


def load_json(path: str | Path) -> Any:
    return json.loads((ROOT / path).read_text(encoding="utf-8"))


def set_wide_layout(prs: Presentation) -> None:
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)


def add_textbox(slide, left: float, top: float, width: float, height: float, text: str, size: int, bold: bool = False, color: RGBColor | None = None) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.05)
    frame.margin_right = Inches(0.05)
    frame.margin_top = Inches(0.03)
    frame.margin_bottom = Inches(0.03)
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.font.name = "Arial"
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    if color is not None:
        paragraph.font.color.rgb = color


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    add_textbox(slide, 0.45, 0.22, 12.4, 0.45, title, 23, bold=True, color=RGBColor(35, 35, 35))
    if subtitle:
        add_textbox(slide, 0.47, 0.72, 12.2, 0.35, subtitle, 10, color=RGBColor(92, 92, 92))


def add_footer(slide, text: str) -> None:
    add_textbox(slide, 0.45, 7.05, 12.4, 0.22, text, 7, color=RGBColor(105, 105, 105))


def add_bullets(slide, left: float, top: float, width: float, height: float, bullets: list[str], size: int = 15) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.name = "Arial"
        p.font.size = Pt(size)


def style_table(table, header_rgb: RGBColor = RGBColor(235, 239, 245), font_size: int = 8) -> None:
    for row_idx, row in enumerate(table.rows):
        for cell in row.cells:
            cell.margin_left = Inches(0.03)
            cell.margin_right = Inches(0.03)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_rgb
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.name = "Arial"
                paragraph.font.size = Pt(font_size)
                paragraph.alignment = PP_ALIGN.CENTER


def add_simple_table(slide, rows: list[list[str]], left: float, top: float, width: float, height: float, font_size: int = 8) -> None:
    table = slide.shapes.add_table(len(rows), len(rows[0]), Inches(left), Inches(top), Inches(width), Inches(height)).table
    for r_idx, row in enumerate(rows):
        for c_idx, value in enumerate(row):
            table.cell(r_idx, c_idx).text = value
    style_table(table, font_size=font_size)


def add_placeholder(slide, left: float, top: float, width: float, height: float, title: str, body: str) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(247, 249, 252)
    shape.line.color.rgb = RGBColor(160, 170, 185)
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Arial"
    p.font.bold = True
    p.font.size = Pt(13)
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = body
    p2.font.name = "Arial"
    p2.font.size = Pt(8)
    p2.alignment = PP_ALIGN.CENTER


def target_rows(summary: dict[str, Any]) -> list[list[str]]:
    names = ["handcrafted", "safe_gaussian", "safe_empirical", "safe_blended"]
    labels = ["Original handcrafted", "Safe Gaussian", "Safe Empirical", "Safe Blended"]
    rows = [["Target", "mu [d_goal,d_hazard,d_agent,speed]", "sigma"]]
    for label, key in zip(labels, names):
        value = summary[key] if key in summary else summary[key.replace("_", " ")]
        rows.append([label, fmt_list(value["mu"]), fmt_list(value["sigma"])])
    return rows


def fmt_list(values: list[float]) -> str:
    return "[" + ", ".join(f"{float(v):.3f}" for v in values) + "]"


def training_rows(rows: list[dict[str, Any]]) -> list[list[str]]:
    table = [["Run", "Return", "Success", "Cost", "StateW2", "FinalW2", "TailRisk", "Unsafe"]]
    for row in rows:
        table.append(
            [
                str(row["name"]).replace("_", " "),
                f"{float(row['return']):.3f}",
                f"{float(row['success']):.2f}",
                f"{float(row['cost']):.1f}",
                f"{float(row['state_w2']):.3f}",
                f"{float(row['final_state_w2']):.3f}",
                f"{float(row['tail_risk']):.3f}",
                f"{float(row['unsafe_occupancy']):.3f}",
            ]
        )
    return table


def distribution_rows(rows: list[dict[str, Any]]) -> list[list[str]]:
    table = [["Run", "Target", "Return", "Cost", "Dist", "TailRisk"]]
    for row in rows:
        table.append(
            [
                str(row["run"]).replace("Method", "M"),
                str(row["target"]).replace("rho_star_", "").replace(".json", ""),
                f"{float(row['mean_return']):.3f}",
                f"{float(row['average_cost']):.1f}",
                f"{float(row['state_distance']):.3f}",
                f"{float(row['tail_risk']):.3f}",
            ]
        )
    return table


def add_paths_slide(prs: Presentation, title: str, paths: list[tuple[str, str]]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, title, "图片不嵌入 PPT，下面列出建议插图和文件位置。")
    y = 1.22
    for label, path in paths:
        add_textbox(slide, 0.65, y, 2.5, 0.27, label, 9, bold=True, color=RGBColor(60, 70, 85))
        add_textbox(slide, 3.1, y, 9.6, 0.27, rel(path), 7, color=RGBColor(70, 70, 70))
        y += 0.42


def build_presentation(output: Path) -> None:
    build_summary = load_json("outputs/targets/safe_rhostar_100k/safe_rhostar_build_summary.json")
    train_summary = load_json("outputs/comparisons/safe_rhostar_training_100k/safe_rhostar_training_comparison.json")
    dist_summary = load_json("outputs/comparisons/safe_rhostar_training_100k/distribution_diagnostics/rho_vs_rhostar_distribution_summary.json")

    prs = Presentation()
    set_wide_layout(prs)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Safe rho* 构造与训练对比总结", "方法三、方法四在 Safe Gaussian / Safe Empirical / Safe Blended 三种 rho* 下的结果")
    add_bullets(
        slide,
        0.85,
        1.45,
        11.8,
        4.7,
        [
            "目标：用高质量 rollout 样本构造更可达的 rho*，并比较方法三、方法四在不同目标分布下的表现。",
            "三种 rho*：Safe Gaussian、Safe Empirical、Safe Blended。",
            "Safe Empirical 使用 full empirical target，保留真实样本分布形状；Gaussian 只保留均值和方差。",
            "训练对比共 6 组：Method3/Method4 x 三种 rho*。",
            "额外画出每组训练后 rollout rho 与对应 rho* 的分布对比图。",
        ],
        size=17,
    )
    add_footer(slide, f"主要结果目录：{rel('outputs/comparisons/safe_rhostar_training_100k')}")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "rho* 构造方式", "从 25 条候选 episode 中筛出高质量 rollout 样本，再构造三种目标分布。")
    add_bullets(
        slide,
        0.65,
        1.18,
        5.7,
        5.5,
        [
            f"候选 episode：{build_summary['candidate_episodes']}",
            f"筛选 episode：{build_summary['selected_episodes']}",
            f"筛选样本数：{build_summary['selected_samples']}",
            "Episode 条件：success=true，return 高，cost/tail risk 低。",
            "Sample 条件：d_hazard >= 0.45，d_agent >= 0.45，speed <= 1.0。",
            "Blended：0.40 * handcrafted + 0.60 * safe_gaussian。",
        ],
        size=14,
    )
    add_simple_table(slide, target_rows(build_summary), 6.55, 1.25, 6.25, 3.9, font_size=7)
    add_footer(slide, f"构造摘要：{rel('outputs/targets/safe_rhostar_100k/safe_rhostar_build_summary.md')}")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "三种 rho* 的直观差异", "这里放 rho* 目标分布对比图。")
    add_placeholder(
        slide,
        0.8,
        1.45,
        11.7,
        4.8,
        "建议插图：safe_rhostar_targets_comparison.png",
        rel("outputs/targets/safe_rhostar_100k/figures/safe_rhostar_targets_comparison.png"),
    )
    add_footer(slide, "这张图用于解释：Gaussian/Empirical 均值方差相同，但 Empirical 保留样本形状。")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "训练对比表", "正式 100k 训练，10 episode evaluation。")
    add_simple_table(slide, training_rows(train_summary), 0.35, 1.16, 12.65, 4.6, font_size=7)
    add_bullets(
        slide,
        0.6,
        6.05,
        12.1,
        0.75,
        [
            "经验版分布距离最低，但 cost/tail risk 明显升高；Gaussian 更安全但更容易保守；Blended 是折中方向。",
        ],
        size=11,
    )
    add_footer(slide, f"对比表：{rel('outputs/comparisons/safe_rhostar_training_100k/safe_rhostar_training_comparison.md')}")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "方法三结果解读", "Method3 对 rho* 类型非常敏感。")
    add_bullets(
        slide,
        0.75,
        1.2,
        12.0,
        5.2,
        [
            "Method3 + Safe Gaussian：Cost 和 TailRisk 最低，但 Return=-5.22、Success=0.30，表现为保守失败。",
            "Method3 + Safe Empirical：StateW2=0.615、FinalW2=0.371，最贴近经验目标；但 Cost=309.4、TailRisk=0.387，安全代价最高。",
            "Method3 + Safe Blended：Return=0.577、Success=0.70、Cost=129.4、TailRisk=0.196，是方法三里较平衡的版本。",
            "结论：如果追求分布贴近，empirical 最强；如果追求综合任务和安全，blended 更合理。",
        ],
        size=15,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "方法四结果解读", "Method4 在新 rho* 下没有超过原始 tail-warmup 主结果。")
    add_bullets(
        slide,
        0.75,
        1.2,
        12.0,
        5.2,
        [
            "Method4 + Safe Gaussian：Return=0.018、Success=0.50、Cost=103.7，偏保守。",
            "Method4 + Safe Empirical：StateW2=0.894、FinalW2=0.612，分布贴近较好；但 Cost=199.4、TailRisk=0.269。",
            "Method4 + Safe Blended：Return=0.938 最高，但 Success=0.40，StateW2=2.018，分布贴近不足。",
            "结论：方法四的 warmup 与新 rho* 需要重新调 penalty；直接替换 rho* 并不能自动提升综合表现。",
        ],
        size=15,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "rho 与 rho* 分布诊断表", "3 episode rollout 的真实分布样本诊断。")
    add_simple_table(slide, distribution_rows(dist_summary), 0.6, 1.15, 12.15, 4.6, font_size=7)
    add_footer(slide, f"诊断表：{rel('outputs/comparisons/safe_rhostar_training_100k/distribution_diagnostics/rho_vs_rhostar_distribution_summary.md')}")

    add_paths_slide(
        prs,
        "建议放入的图：目标 rho*",
        [
            ("三种 rho* 对比", "outputs/targets/safe_rhostar_100k/figures/safe_rhostar_targets_comparison.png"),
            ("rho* 构造摘要", "outputs/targets/safe_rhostar_100k/safe_rhostar_build_summary.md"),
        ],
    )

    add_paths_slide(
        prs,
        "建议放入的图：Method3",
        [
            ("M3 Gaussian rho 对比", "outputs/comparisons/safe_rhostar_training_100k/distribution_diagnostics/method3_safe_gaussian_rho_vs_rhostar.png"),
            ("M3 Empirical rho 对比", "outputs/comparisons/safe_rhostar_training_100k/distribution_diagnostics/method3_safe_empirical_rho_vs_rhostar.png"),
            ("M3 Blended rho 对比", "outputs/comparisons/safe_rhostar_training_100k/distribution_diagnostics/method3_safe_blended_rho_vs_rhostar.png"),
        ],
    )

    add_paths_slide(
        prs,
        "建议放入的图：Method4",
        [
            ("M4 Gaussian rho 对比", "outputs/comparisons/safe_rhostar_training_100k/distribution_diagnostics/method4_safe_gaussian_rho_vs_rhostar.png"),
            ("M4 Empirical rho 对比", "outputs/comparisons/safe_rhostar_training_100k/distribution_diagnostics/method4_safe_empirical_rho_vs_rhostar.png"),
            ("M4 Blended rho 对比", "outputs/comparisons/safe_rhostar_training_100k/distribution_diagnostics/method4_safe_blended_rho_vs_rhostar.png"),
        ],
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "结论与下一步", "当前新 rho* 实验说明目标分布选择会显著改变任务/安全折中。")
    add_bullets(
        slide,
        0.75,
        1.18,
        12.0,
        5.5,
        [
            "Safe Gaussian：稳定、快、偏保守；容易降低 cost，但可能远离任务目标。",
            "Safe Empirical：最贴近经验分布；但可能复现高质量样本里的风险模式，cost/tail risk 偏高。",
            "Safe Blended：当前最适合作为后续调参起点，尤其是 Method3 + Blended。",
            "下一步建议：固定 blended rho*，调小/调度 distribution penalty，配合 cost/tail warmup，重新训练 Method4。",
            "论文表述上：rho* 不应只靠手工设定，也不应直接复制成功轨迹；需要用安全筛选后的可达目标分布。",
        ],
        size=15,
    )

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--output", default="outputs/presentations/safe_rhostar_experiment_summary.pptx")
    args = parser.parse_args()
    build_presentation(ROOT / args.output)
    print(f"saved presentation to {ROOT / args.output}")


if __name__ == "__main__":
    main()
