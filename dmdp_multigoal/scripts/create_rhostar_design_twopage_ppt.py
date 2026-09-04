#!/usr/bin/env python
from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]


def load_json(path: str | Path) -> Any:
    return json.loads((ROOT / path).read_text(encoding="utf-8"))


def set_wide_layout(prs: Presentation) -> None:
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)


def add_text(slide, left: float, top: float, width: float, height: float, text: str, size: int, bold: bool = False, color: RGBColor | None = None) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Arial"
    p.font.size = Pt(size)
    p.font.bold = bold
    if color is not None:
        p.font.color.rgb = color


def style_table(table, font_size: int = 8) -> None:
    for r, row in enumerate(table.rows):
        for cell in row.cells:
            cell.margin_left = Inches(0.03)
            cell.margin_right = Inches(0.03)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(232, 237, 245)
            for p in cell.text_frame.paragraphs:
                p.font.name = "Arial"
                p.font.size = Pt(font_size)
                p.alignment = PP_ALIGN.CENTER


def add_table(slide, rows: list[list[str]], left: float, top: float, width: float, height: float, font_size: int = 8) -> None:
    table = slide.shapes.add_table(len(rows), len(rows[0]), Inches(left), Inches(top), Inches(width), Inches(height)).table
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            table.cell(r, c).text = value
    style_table(table, font_size)


def fmt(values: list[float]) -> str:
    return "[" + ", ".join(f"{float(v):.3f}" for v in values) + "]"


def title(slide, text: str, subtitle: str) -> None:
    add_text(slide, 0.45, 0.18, 12.4, 0.42, text, 22, True, RGBColor(35, 35, 35))
    add_text(slide, 0.47, 0.63, 12.1, 0.28, subtitle, 10, False, RGBColor(90, 90, 90))


def build(output: Path) -> None:
    build_summary = load_json("outputs/targets/safe_rhostar_100k/safe_rhostar_build_summary.json")
    train_rows = load_json("outputs/comparisons/safe_rhostar_training_100k/safe_rhostar_training_comparison.json")

    prs = Presentation()
    set_wide_layout(prs)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title(slide, "三种 rho* 的设计", "rho* 定义在 x=[d_goal,d_hazard,d_agent,speed] 的状态特征分布上")
    bullets = [
        "真实最优安全分布 rho_{pi_safe*} 不可直接观测，因此用高质量 rollout 样本构造近似目标。",
        f"筛选：{build_summary['candidate_episodes']} 条候选 episode -> {build_summary['selected_episodes']} 条高质量 episode -> {build_summary['selected_samples']} 个样本。",
        "Episode 条件：success=true，return 高，cost/tail risk 低；Sample 条件：d_hazard>=0.45，d_agent>=0.45，speed<=1.0。",
    ]
    y = 1.05
    for b in bullets:
        add_text(slide, 0.65, y, 12.1, 0.35, "• " + b, 13)
        y += 0.42
    target_table = [
        ["Target", "形式", "mu", "sigma", "含义"],
        ["Safe Gaussian", "N(mu_good,diag(sigma_good^2))", fmt(build_summary["safe_gaussian"]["mu"]), fmt(build_summary["safe_gaussian"]["sigma"]), "可达高质量样本的高斯压缩"],
        ["Safe Empirical", "Empirical({x_good})", fmt(build_summary["safe_empirical"]["mu"]), fmt(build_summary["safe_empirical"]["sigma"]), "保留真实样本形状"],
        ["Safe Blended", "0.4 handcrafted + 0.6 gaussian", fmt(build_summary["safe_blended"]["mu"]), fmt(build_summary["safe_blended"]["sigma"]), "理论安全与经验可达折中"],
    ]
    add_table(slide, target_table, 0.35, 2.55, 12.65, 2.0, 7)
    diff_table = [
        ["rho*", "优点", "主要风险", "适合用途"],
        ["Gaussian", "快、稳定、易优化", "过度简化，可能保守", "可达高斯目标实验"],
        ["Empirical", "最真实，能表达非高斯", "慢，可能继承风险行为", "检验真实经验分布匹配"],
        ["Blended", "任务/安全折中", "权重需要调参", "主方法后续调参起点"],
    ]
    add_table(slide, diff_table, 0.65, 5.0, 12.0, 1.45, 8)
    add_text(slide, 0.55, 6.82, 12.2, 0.25, "建议插图：outputs/targets/safe_rhostar_100k/figures/safe_rhostar_targets_comparison.png", 8, False, RGBColor(95, 95, 95))

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title(slide, "三种 rho* 的训练结果与判断", "Method3/Method4 分别使用三种 rho* 训练 100k，并做 10 episode evaluation")
    rows = [["Run", "Return", "Success", "Cost", "StateW2", "FinalW2", "TailRisk", "Unsafe"]]
    for row in train_rows:
        rows.append([
            str(row["name"]).replace("_", " "),
            f"{float(row['return']):.3f}",
            f"{float(row['success']):.2f}",
            f"{float(row['cost']):.1f}",
            f"{float(row['state_w2']):.3f}",
            f"{float(row['final_state_w2']):.3f}",
            f"{float(row['tail_risk']):.3f}",
            f"{float(row['unsafe_occupancy']):.3f}",
        ])
    add_table(slide, rows, 0.35, 1.12, 12.65, 3.0, 7)
    conclusions = [
        "Safe Empirical：分布距离最低，但 cost/tail risk 偏高，说明它会强力贴近经验分布，也可能复制经验样本中的风险模式。",
        "Safe Gaussian：更安全但容易保守，Method3 中 return 和 success 明显下降。",
        "Safe Blended：目前更适合作为后续主实验起点，尤其 Method3 + Blended 在任务与安全之间更平衡。",
        "后续建议：固定 Safe Blended，重新调 distribution/cost/tail penalty 和 warmup，而不是直接替换 rho* 后不调参数。",
    ]
    y = 4.45
    for b in conclusions:
        add_text(slide, 0.7, y, 12.0, 0.38, "• " + b, 12)
        y += 0.48
    add_text(slide, 0.55, 6.55, 12.2, 0.22, "建议插图：outputs/comparisons/safe_rhostar_training_100k/distribution_diagnostics/summary_plots/rho_distance_return_cost_tradeoff.png", 8, False, RGBColor(95, 95, 95))
    add_text(slide, 0.55, 6.82, 12.2, 0.22, "建议插图：outputs/comparisons/safe_rhostar_training_100k/distribution_diagnostics/summary_plots/rho_rhostar_gap_heatmaps.png", 8, False, RGBColor(95, 95, 95))

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--output", default="outputs/presentations/rhostar_design_twopage.pptx")
    args = parser.parse_args()
    output = ROOT / args.output
    build(output)
    print(f"saved presentation to {output}")


if __name__ == "__main__":
    main()
