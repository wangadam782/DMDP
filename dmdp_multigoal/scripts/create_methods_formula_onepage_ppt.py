#!/usr/bin/env python
from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


METHODS = [
    {
        "name": "Method 1\nState-MAPPO",
        "policy": "a_t^i ~ pi_theta(a_t^i | o_t^i)",
        "reward": "r_train = r_env - lambda_c c_t",
        "critic": "V_phi(s_t)",
        "note": "baseline: no rho feedback",
        "color": RGBColor(76, 120, 168),
    },
    {
        "name": "Method 2\nDist-MAPPO",
        "policy": "a_t^i ~ pi_theta(a_t^i | o_t^i)",
        "reward": "r_train = r_env - lambda_c c_t",
        "critic": "Z_phi(s_t) = {z_k(s_t)}",
        "note": "return distribution critic",
        "color": RGBColor(114, 183, 178),
    },
    {
        "name": "Method 2-rho\nDist-MAPPO + rho",
        "policy": "omega_t=[mu_t,sigma_t]\na_t^i ~ pi_theta(a_t^i | o_t^i, omega_t)",
        "reward": "r_train = r_env - lambda_c c_t",
        "critic": "Z_phi(s_t, omega_t)",
        "note": "same reward, adds rho feedback",
        "color": RGBColor(84, 162, 75),
    },
    {
        "name": "Method 3\nDMDP-MAPPO",
        "policy": "omega_t=[mu_t,sigma_t]\na_t^i ~ pi_theta(a_t^i | o_t^i, omega_t)",
        "reward": "r_train = r_env - lambda_c c_t\n          - lambda_d W2(rho_t,rho*)\n          - lambda_tail R_tail",
        "critic": "V_phi(s_t, omega_t)",
        "note": "engineering DMDP",
        "color": RGBColor(229, 86, 86),
    },
    {
        "name": "Method 4\nPaper-like Online",
        "policy": "Delta omega_t=omega_t-omega_{t-1}\nu_t=phi_eta(omega_t,Delta omega_t)\na_t^i ~ pi_theta(a_t^i | o_t^i,u_t)",
        "reward": "r_train = r_env - lambda_c(t)c_t\n          - lambda_d W2(rho_t,rho*)\n          - lambda_tail(t)R_tail\n          - lambda_L(t)L_violation",
        "critic": "V_phi(omega_t,Delta omega_t)\nL_psi(omega_t,Delta omega_t)",
        "note": "main theory version",
        "color": RGBColor(245, 133, 24),
    },
    {
        "name": "Method 5\nOnline + empirical rho*",
        "policy": "u_t=phi_eta(omega_t,Delta omega_t)\na_t^i ~ pi_theta(a_t^i | o_t^i,u_t)",
        "reward": "r_train = r_env - lambda_c(t)c_t\n          - lambda_d D_SW(rho_t,rho*_M2)\n          - lambda_tail(t)R_tail\n          - lambda_L(t)L_violation",
        "critic": "V_phi(omega_t,Delta omega_t)\nL_psi(omega_t,Delta omega_t)",
        "note": "rho*_M2 from Method 2 rollouts",
        "color": RGBColor(178, 121, 162),
    },
]


def set_wide_layout(prs: Presentation) -> None:
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)


def add_textbox(slide, left: float, top: float, width: float, height: float, text: str, size: int, bold: bool = False, color: RGBColor | None = None) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.04)
    frame.margin_right = Inches(0.04)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.name = "Arial"
    if color is not None:
        paragraph.font.color.rgb = color


def add_method_card(slide, method: dict, left: float, top: float, width: float, height: float) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(248, 249, 251)
    shape.line.color.rgb = method["color"]
    shape.line.width = Pt(2.0)
    shape.text_frame.clear()

    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(0.12),
        Inches(height),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = method["color"]
    bar.line.fill.background()

    add_textbox(slide, left + 0.22, top + 0.10, width - 0.34, 0.38, method["name"], 10, bold=True, color=method["color"])
    add_textbox(slide, left + 0.22, top + 0.52, width - 0.34, 0.34, "Policy feedback", 7, bold=True, color=RGBColor(80, 80, 80))
    add_textbox(slide, left + 0.22, top + 0.78, width - 0.34, 0.58, method["policy"], 7)
    add_textbox(slide, left + 0.22, top + 1.38, width - 0.34, 0.26, "Reward", 7, bold=True, color=RGBColor(80, 80, 80))
    add_textbox(slide, left + 0.22, top + 1.62, width - 0.34, 0.74, method["reward"], 6)
    add_textbox(slide, left + 0.22, top + 2.38, width - 0.34, 0.26, "Critic", 7, bold=True, color=RGBColor(80, 80, 80))
    add_textbox(slide, left + 0.22, top + 2.62, width - 0.34, 0.45, method["critic"], 6)
    add_textbox(slide, left + 0.22, top + height - 0.35, width - 0.34, 0.22, method["note"], 6, color=RGBColor(95, 95, 95))


def add_legend(slide) -> None:
    box = slide.shapes.add_textbox(Inches(0.35), Inches(6.95), Inches(12.6), Inches(0.28))
    frame = box.text_frame
    frame.clear()
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.text = "omega_t parameterizes rho_t; rho* is the target state distribution; D_SW is sliced-Wasserstein distance to the Method 2 empirical rho*."
    paragraph.font.size = Pt(8)
    paragraph.font.color.rgb = RGBColor(90, 90, 90)


def build_presentation(output: Path) -> None:
    prs = Presentation()
    set_wide_layout(prs)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_textbox(slide, 0.35, 0.16, 12.7, 0.42, "方法对比：策略反馈、Reward、Critic 的公式差异", 21, bold=True, color=RGBColor(35, 35, 35))
    add_textbox(
        slide,
        0.36,
        0.62,
        12.4,
        0.28,
        "从普通状态反馈，到 rho_t 分布反馈，再到方法四的高层在线分布反馈控制。",
        10,
        color=RGBColor(90, 90, 90),
    )

    card_w = 4.05
    card_h = 2.85
    x_positions = [0.35, 4.64, 8.93]
    y_positions = [1.05, 4.03]
    for idx, method in enumerate(METHODS):
        row = idx // 3
        col = idx % 3
        add_method_card(slide, method, x_positions[col], y_positions[row], card_w, card_h)

    add_legend(slide)
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--output", default="outputs/presentations/method_formula_comparison_onepage.pptx")
    args = parser.parse_args()
    build_presentation(Path(args.output))
    print(f"saved presentation to {args.output}")


if __name__ == "__main__":
    main()
